#!/usr/bin/env python3
"""Genere la presentation de soutenance (INF462) au format PowerPoint."""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLEU = RGBColor(0x1E, 0x3A, 0x8A)
BLEU2 = RGBColor(0x25, 0x63, 0xEB)
GRIS = RGBColor(0x33, 0x41, 0x55)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def fond(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def zone(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def slide_titre(titre, sous_titre):
    s = prs.slides.add_slide(BLANK)
    fond(s, BLEU)
    tf = zone(s, 0.8, 2.4, 11.7, 2.5)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = titre
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BLANC
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = sous_titre
    r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    return s


def slide_contenu(titre, points):
    s = prs.slides.add_slide(BLANK)
    fond(s, BLANC)
    # bandeau titre
    bar = s.shapes.add_textbox(Inches(0), Inches(0.3), Inches(13.333), Inches(0.9))
    tf = bar.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = titre
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BLEU
    # contenu
    tf = zone(s, 0.9, 1.5, 11.6, 5.6)
    first = True
    for niveau, texte in points:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = niveau
        r = p.add_run(); r.text = ("• " if niveau == 0 else "– ") + texte
        r.font.size = Pt(20 if niveau == 0 else 16)
        r.font.color.rgb = GRIS if niveau == 0 else RGBColor(0x47, 0x55, 0x69)
        if niveau == 0:
            r.font.bold = True
        p.space_after = Pt(6)
    return s


# ---- Slides ----
slide_titre("Plateforme Bancaire Distribuée",
            "INF462 — Architecture Logicielle · Microservices · Session 2025/2026")

slide_contenu("Contexte & Objectif", [
    (0, "Plateforme financière distribuée pour plusieurs opérateurs (banques, microfinances, mobile)"),
    (1, "Volume croissant de transactions, disponibilité, sécurité, traçabilité, scalabilité"),
    (0, "Expérience homogène, sécurisée et transparente pour tous les utilisateurs"),
    (0, "Démarche complète : analyse métier (DDD) → développement → déploiement → supervision"),
])

slide_contenu("Cahier des charges (extrait)", [
    (0, "Exigences fonctionnelles"),
    (1, "Clients & KYC, comptes, dépôts/retraits, transferts intra & inter-opérateurs"),
    (1, "Prêts : demande, validation, échéancier, remboursement"),
    (1, "Notifications, opérateurs, rapports, audit, OCR de documents"),
    (1, "Plusieurs mécanismes d'authentification"),
    (0, "Exigences non fonctionnelles"),
    (1, "Disponibilité, scalabilité, sécurité, résilience, observabilité, évolutivité"),
])

slide_contenu("Analyse Domain-Driven Design (DDD)", [
    (0, "Sous-domaines"),
    (1, "Core : Prêts, Transactions  ·  Supporting : Comptes, Clients, Documents  ·  Generic : Auth, Notifications"),
    (0, "Bounded Contexts → 1 microservice = 1 base (database per service)"),
    (1, "Identity, Customer, Account, Transaction, Loan, Document/AI, Notification"),
    (0, "Communication : par identifiant + événements (pas de jointure inter-bases)"),
])

slide_contenu("Architecture microservices", [
    (0, "Frontend Angular (Nginx) → API Gateway (point d'entrée unique)"),
    (0, "Socle Cloud Native"),
    (1, "API Gateway (routage + sécurité), Eureka (découverte), Config Server (config centralisée)"),
    (0, "7 microservices métier + infrastructure (PostgreSQL, RabbitMQ)"),
    (0, "Architecture polyglotte imposée : Java + JavaScript + Python"),
])

slide_contenu("Stack technologique (polyglotte)", [
    (0, "Java / Spring Boot : auth, customer, account, transaction, loan, gateway, config, discovery"),
    (0, "Python / FastAPI : ai-document (OCR Tesseract)"),
    (0, "Node.js : notification (consommateur RabbitMQ + email)"),
    (0, "Angular : interface utilisateur (espaces client / opérateur / admin)"),
    (0, "PostgreSQL (1 base par service) · RabbitMQ (événements) · Docker"),
])

slide_contenu("Sécurité", [
    (0, "Authentification — 2 mécanismes"),
    (1, "Email + mot de passe (haché BCrypt) → jeton JWT"),
    (1, "Connexion Google (vérification de l'ID token, OAuth2)"),
    (0, "JWT stateless vérifié à la Gateway (filtre global) avant routage"),
    (0, "Propagation de l'identité (X-User-Email / X-User-Roles) + autorisation par rôle"),
])

slide_contenu("Communications synchrones & asynchrones", [
    (0, "Synchrone (REST) via la Gateway + Eureka (lb://)"),
    (1, "Ex. transaction-service → account-service (crédit/débit)"),
    (0, "Asynchrone via RabbitMQ (exchange topic banking.events)"),
    (1, "transaction.completed → notification-service (log + email)"),
    (0, "Résilience : circuit breaker (resilience4j) sur les appels inter-services"),
])

slide_contenu("Gestion documentaire & IA (OCR)", [
    (0, "Soumission de documents (CNI, passeport, justificatifs…)"),
    (0, "Extraction automatique du texte par OCR (Tesseract, service Python distribué)"),
    (0, "Alimente le processus métier : validation du KYC du client"),
    (0, "Service distribué, intégré à l'architecture via la Gateway"),
])

slide_contenu("Cloud Native & DevOps", [
    (0, "Conteneurisation complète : Docker + docker-compose (14 conteneurs)"),
    (0, "Orchestration : manifests Kubernetes (namespace, services, ingress)"),
    (0, "CI/CD : GitHub Actions (build + tests Java/Python/Node/Angular + images)"),
    (0, "Observabilité : Prometheus (métriques) + Grafana (dashboards)"),
    (0, "Traçabilité : audit centralisé des requêtes à la Gateway"),
])

slide_contenu("Démonstration", [
    (0, "Connexion admin → espaces différenciés par rôle"),
    (0, "Créer opérateur → client (KYC) → comptes"),
    (0, "Dépôt, retrait, transfert (soldes mis à jour, commission)"),
    (0, "Demande de prêt → approbation → échéancier → remboursement"),
    (0, "OCR d'un document → validation KYC"),
    (0, "Notification temps réel (asynchrone) + métriques Grafana"),
])

slide_contenu("Difficultés & compromis", [
    (0, "Hétérogénéité des contributions (variables d'env, brokers) → harmonisation"),
    (0, "Spring Cloud 2025.x : renommage gateway + transport Eureka / load-balancer"),
    (0, "Builds Maven en conteneur instables → images jar-préconstruit"),
    (0, "Compromis : cohérence des transferts (cas simple + circuit breaker, saga en perspective)"),
])

slide_contenu("Perspectives & Conclusion", [
    (0, "Saga / compensation pour les transferts inter-opérateurs"),
    (0, "Vérification auto OCR + scoring de prêt par IA"),
    (0, "Tracing distribué (OpenTelemetry), audit persistant, déploiement K8s + autoscaling"),
    (0, "Plateforme microservices complète, sécurisée, observable et conforme aux exigences"),
])

slide_titre("Merci — Questions ?", "Plateforme Bancaire Distribuée · INF462 · 2025/2026")

import os
out = os.path.join(os.path.dirname(__file__), "INF462-presentation.pptx")
prs.save(out)
print("✅ Présentation générée :", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
